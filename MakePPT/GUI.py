# -*- coding:utf-8 -*-
from pptx import Presentation
from PyQt5.QtWidgets import (QApplication,QFileDialog, QComboBox, QVBoxLayout,QWidget, QPushButton, QDesktopWidget,QGridLayout, QLabel, QLineEdit,QRadioButton)
from PyQt5.QtGui import QIcon, QFont
from PyQt5.QtCore import QCoreApplication
import os
import sys
import io
import copy

sys.stdout = io.TextIOWrapper(sys.stdout.detach(), encoding = 'utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.detach(), encoding = 'utf-8')

old_array = ['창세기','출애굽기','레위기','민수기','신명기','여호수아','사사기','룻기','사무엘상','사무엘하','열왕기상','열왕기하','역대상','역대하','에스라','느혜미야',
               '에스더','욥기','시편','잠언','전도서','아가','이사야','예레미야','예레미야애가','에스겔','다니엘','호세아','요엘','아모스','오바댜','요나','미가','나훔',
               '하박국','스바냐','학개','스가랴','말라기']
new_array = ['마태복음','마가복음','누가복음','요한복음','사도행전','로마서','고린도전서','고린도후서','갈라디아서','에베소서','빌립보서','골로새서','데살로니가전서','데살로니가후서',
             '디모데전서','디모데후서','디도서','빌레몬서','히브리서','야고보서','베드로전서','베드로후서','요한일서','요한이서','요한삼서','유다서','요한계시록']

class MyApp(QWidget):

    def __init__(self):
        super().__init__()
        self.initUI()

    def initUI(self):

        self.lbl = QLabel('말씀 선택',self)
        self.lbl.move(0,0)
        cb = QComboBox(self)
        cb.addItems(old_array)

        chapter = QLineEdit(self)
        start = QLineEdit(self)
        end = QLineEdit(self)
        title = QLineEdit(self)
        old = QRadioButton('구약', self)
        new = QRadioButton('신약', self)
        pptPath = QLineEdit()
        txtPath = QLineEdit()
        fileSearch = QPushButton("PPT 파일 찾기",self)
        wordsFileSearch = QPushButton("말씀 폴더 찾기",self)

        old.setChecked(True)

        grid = QGridLayout()
        self.setLayout(grid)

        grid.addWidget(QLabel('말씀:'), 0, 0)
        grid.addWidget(QLabel('장:'), 1, 0)
        grid.addWidget(QLabel('시작 절:'), 2, 0)
        grid.addWidget(QLabel('끝 절:'), 3, 0)
        grid.addWidget(QLabel('제목:'), 4, 0)
        grid.addWidget(QLabel('PPT 양식:'),5,0)
        grid.addWidget(QLabel('말씀 폴더:'),6,0)

        old.move(60,20)
        new.move(120,20)

        grid.addWidget(cb, 0, 1)
        grid.addWidget(chapter, 1, 1)
        grid.addWidget(start, 2, 1)
        grid.addWidget(end, 3, 1)
        grid.addWidget(title, 4, 1)
        grid.addWidget(pptPath,5,1)
        grid.addWidget(fileSearch,5,2)
        grid.addWidget(txtPath,6,1)
        grid.addWidget(wordsFileSearch,6,2)

        cb.activated[str].connect(self.onActivated)

        btn = QPushButton('Start', self)
        btn.move(400,440)

        self.setWindowTitle('making ppt Automation program - Omega Church')
        self.setWindowIcon(QIcon('omega_logo.png'))
        self.resize(900, 490) #사이즈 조절
        self.center() #센터 맞춘다.

        #이벤트 처리 설정
        btn.clicked.connect(lambda: self.make(str(title.text()),str(cb.currentText()),str(chapter.text()),str(start.text()),str(end.text()),str(pptPath.text()),str(txtPath.text()) ) )
        old.clicked.connect(lambda: self.radioButtonClickedOld(cb))
        new.clicked.connect(lambda: self.radioButtonClickedNew(cb))
        fileSearch.clicked.connect(lambda: self.SearchButtonClicked(pptPath))
        wordsFileSearch.clicked.connect(lambda: self.searchWordFileButtonClicked(txtPath))

        self.show() #창을 띄운다.


    def radioButtonClickedOld(self,combobox):
        combobox.clear()
        combobox.addItems(old_array)

    def radioButtonClickedNew(self,combobox):
        combobox.clear()
        combobox.addItems(new_array)

    def searchWordFileButtonClicked(self,path):
        fname=QFileDialog.getExistingDirectory(self)
        path.setText(fname)

    def SearchButtonClicked(self,path):
        fname = QFileDialog.getOpenFileName(self)
        path.setText(fname[0])

    def onActivated(self, text):
        self.lbl.setText(text)
        self.lbl.adjustSize()

    def center(self): #창을 가운데 띄우기위한 함수
        qr = self.frameGeometry()
        cp = QDesktopWidget().availableGeometry().center()
        qr.moveCenter(cp)
        self.move(qr.topLeft())


    def make(self,input_title, input_chapter, input_number, input_start, input_end, input_pptPath, input_txtPath):
        mytitle = input_title
        Chapter = input_chapter
        number = int(input_number)
        start = int(input_start)
        end = int(input_end)
        pPath = input_pptPath
        tPath = input_txtPath

        print("Title : " + Chapter)
        print("<<< Start >>>")
        print()

        pptPath = input_pptPath
        txtName = Chapter + str(object=number) + "장"
        txtPath = tPath+'/'+Chapter+'/'+ txtName+'.txt'
        txtPath = os.path.abspath(txtPath)

        words = []

        try:
            targetFile = open(txtPath, 'rt', encoding='UTF-8')
            lines = targetFile.readlines()

            for line in lines:
                words.append(line.strip())  # word 배열에 한줄씩 저장
            targetFile.close()
        except IOError as err:
            print("Server File Error: " + str(err))
            # 여기까지 왔으면 Words에 말씀 저장완료


        try: #디렉터리 생성
            if not(os.path.isdir('output folder')):
                os.makedirs(os.path.join('output folder'))
        except OSError as e:
            if e.errno != errno.EEXIST:
                print("Failed to create directory!!!!!")
                raise


        # 인자로 불러올 파일 경로를 넣어줄 수 있다. 인자 없으면 현재 경로에 새롭게 생성
        prs = Presentation(pptPath)
        title_slide_layout = prs.slide_layouts[0]  # 제목 슬라이드 레이아웃 지정 Layout 0번이다.
        count=1
        for contents in words:
            (verse, word) = contents.split('.', 1)
            if count >=start and count <=end:
                # 그 레이아웃으로 슬라이드 하나 추가하고 그걸 slide변수라고 지정
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title  # 제목 상자 지정
                subtitle = slide.placeholders[1]  # 부제목 상자 지정
                title.text = str(mytitle)  # 제목 텍스트 적는다.
                subtitle.text = str(verse) + "\n" + str(word.strip())  # 부제목 텍스트 적는다.
                count += 1
            else:
                count += 1
        prs.save('output folder' + './' + str(mytitle) + '.pptx')


if __name__ == '__main__':

    app = QApplication(sys.argv) #어플리케이션 객체 생성
    ex = MyApp()
    #ex.show()
    sys.exit(app.exec_())
