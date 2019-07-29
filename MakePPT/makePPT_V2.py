# -*- coding:utf-8 -*-
from pptx import Presentation
import sys
import io
import os
import copy

sys.stdout = io.TextIOWrapper(sys.stdout.detach(), encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.detach(), encoding='utf-8')

def make(input_title, input_chapter, input_number, input_start, input_end):
    mytitle = input_title
    Chapter = input_chapter
    number = input_number
    start = input_start
    end = input_end

    '''
    장 마다 바꿔주어야 할 값은 위의 두 값밖에 없다. 대신 값들은 정확해야한다.
    Chapter = 장 이름
    howMany = 장 수
    '''
    print("Title : " + Chapter)
    print("<<< Start >>>")
    print()

    txtName = Chapter + str(object=number) + "장"
    pptPath = "D:/minwoo/my_study/ppt_study/test.pptx"

    txtPath = txtName+'.txt'
    txtPath = os.path.abspath(txtPath)

    print(txtPath)

    name = str(copy.deepcopy(txtPath).rsplit('\\', 1)[1].split('.')[0])
    # 이름 parsing

    print(name)
    words = []



    try:
        targetFile = open(txtPath, 'rt', encoding='UTF-8')
        lines = targetFile.readlines()

        for line in lines:
            words.append(line.strip())  # word 배열에 한줄씩 저장
        targetFile.close()
    except IOError as err:
        print("Server File Error: " + str(err))
        # 여기까지 왔으면 Words에 말씀 저장완료



    try: #디렉터리 생성
        if not(os.path.isdir(Chapter)):
            os.makedirs(os.path.join(Chapter))
    except OSError as e:
        if e.errno != errno.EEXIST:
            print("Failed to create directory!!!!!")
            raise

    print('???????????????')



    # 인자로 불러올 파일 경로를 넣어줄 수 있다. 인자 없으면 현재 경로에 새롭게 생성
    prs = Presentation(pptPath)
    #prs = Presentation(pptPath)
    title_slide_layout = prs.slide_layouts[0]  # 제목 슬라이드 레이아웃 지정 Layout 0번이다.
    count=1
    for contents in words:
        (verse, word) = contents.split('.', 1)
        if count >=start and count <=end:
            # 그 레이아웃으로 슬라이드 하나 추가하고 그걸 slide변수라고 지정
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title  # 제목 상자 지정
            subtitle = slide.placeholders[1]  # 부제목 상자 지정
            title.text = str(mytitle)  # 제목 텍스트 적는다.
            subtitle.text = str(verse) + "\n" + str(word.strip())  # 부제목 텍스트 적는다.
            count += 1
        else:
            count += 1

    prs.save(Chapter + './' + str(name) + '.pptx')

#make('내가만든제목','창세기',2,5,7)
#print()
#print("<<< normal exit >>>")
